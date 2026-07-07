"""Tests for office_ingest (pptx/xlsx/docx → sibling markdown pre-render).

Needs python-pptx / openpyxl / python-docx — the same deps the box image bakes
in. Run: python test_office_ingest.py
"""

import tempfile
from pathlib import Path

import office_ingest


def _make_samples(root: Path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "GPU 选型"
    s.placeholders[1].text = "H100 vs A100"
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "拓扑"
    tb = s2.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(5), Inches(1)).table
    tb.cell(0, 0).text = "区域"; tb.cell(0, 1).text = "节点"
    tb.cell(1, 0).text = "华东"; tb.cell(1, 1).text = "52"
    prs.save(str(root / "deck.pptx"))

    from openpyxl import Workbook
    wb = Workbook(); ws = wb.active; ws.title = "配额"
    ws.append(["团队", "GPU"]); ws.append(["train", "32"])
    wb.save(str(root / "sub" / "quota.xlsx"))  # nested → exercises rglob + rel path

    from docx import Document
    d = Document(); d.add_heading("手册", level=1); d.add_paragraph("正文一句。")
    t = d.add_table(rows=1, cols=2)
    t.rows[0].cells[0].text = "现象"; t.rows[0].cells[1].text = "动作"
    d.save(str(root / "manual.docx"))


def test_convert_tree():
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        (root / "sub").mkdir()
        (root / "notes.md").write_text("already markdown", "utf-8")  # non-office → untouched
        _make_samples(root)
        converted, errors = office_ingest.convert_tree(str(root))
        assert not errors, errors
        assert dict(converted) == {
            "deck.pptx": "deck.pptx.md",
            "manual.docx": "manual.docx.md",
            "sub/quota.xlsx": "sub/quota.xlsx.md",
        }, converted
        pptx_md = (root / "deck.pptx.md").read_text("utf-8")
        assert "## Slide 1" in pptx_md and "GPU 选型" in pptx_md
        assert "| 区域 | 节点 |" in pptx_md and "| 华东 | 52 |" in pptx_md  # table preserved
        xlsx_md = (root / "sub" / "quota.xlsx.md").read_text("utf-8")
        assert "## Sheet: 配额" in xlsx_md and "| train | 32 |" in xlsx_md
        docx_md = (root / "manual.docx.md").read_text("utf-8")
        assert docx_md.startswith("# 手册") and "正文一句。" in docx_md and "| 现象 | 动作 |" in docx_md
        assert (root / "notes.md").read_text("utf-8") == "already markdown"  # untouched
        # idempotent: a re-run (same bundle re-installed) converts nothing new
        again, _ = office_ingest.convert_tree(str(root))
        assert again == [], again
    print("OK  convert_tree (pptx/xlsx/docx → sibling md; tables, nested, non-office untouched, idempotent)")


def test_fail_open_on_corrupt():
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        (root / "broken.pptx").write_bytes(b"not a real pptx")  # corrupt
        from docx import Document
        doc = Document(); doc.add_paragraph("fine"); doc.save(str(root / "ok.docx"))
        converted, errors = office_ingest.convert_tree(str(root))
        assert [s for s, _ in converted] == ["ok.docx"], converted   # valid one still rendered
        assert [s for s, _ in errors] == ["broken.pptx"], errors     # corrupt recorded, NOT raised
        assert (root / "ok.docx.md").exists() and not (root / "broken.pptx.md").exists()
    print("OK  fail-open (corrupt file recorded in errors, valid files still converted, never raises)")


def main():
    test_convert_tree()
    test_fail_open_on_corrupt()
    print("ALL OK  test_office_ingest")


if __name__ == "__main__":
    main()
