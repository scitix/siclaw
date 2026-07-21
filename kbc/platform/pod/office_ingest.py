"""office_ingest — pre-render binary office sources to markdown the agent can Read.

The served box receives raw/ as ORIGINAL bytes. The agent's Read tool (and the
model) handle pdf / images / plain text natively, but zip-based office formats
(.pptx / .xlsx / .docx) are opaque binary. At raw-install time this renders each
into a sibling `<name>.md` (e.g. deck.pptx -> deck.pptx.md) so the agent reads
clean markdown; the ORIGINAL file is left in place for provenance.

Deliberately lightweight per-format extraction (python-pptx / openpyxl /
python-docx) — no OCR / vision. Images inside a deck stay native for the
compile-stage model. Imports are lazy so this module stays importable without
the optional deps (e.g. a dep-less unit host); an actual conversion needs them
installed — they are baked into the box image.
"""
from __future__ import annotations

import os
import tempfile
import zipfile
from collections.abc import Iterable, Iterator
from pathlib import Path

OFFICE_EXTS = (".pptx", ".xlsx", ".docx")
DEFAULT_MAX_DERIVED_BYTES = 512 * 1024 * 1024
DEFAULT_MAX_ARCHIVE_UNPACKED_BYTES = 512 * 1024 * 1024
DEFAULT_MAX_ARCHIVE_FILES = 10_000


class OfficeIngestLimitExceeded(ValueError):
    """An Office source would exceed a declared compile-box resource budget."""


def _positive_env(name: str, default: int) -> int:
    value = int(os.environ.get(name, str(default)))
    if value <= 0:
        raise ValueError(f"{name} must be positive")
    return value


def _table_lines(rows: Iterable[list[str]]) -> Iterator[str]:
    """Stream a row matrix as a GitHub markdown table (first row = header)."""
    iterator = iter(rows)
    try:
        header = next(iterator)
    except StopIteration:
        return
    width = len(header)

    def cells(r: list[str]) -> str:
        r = (r + [""] * width)[:width]  # pad/truncate to header width → valid table
        return "| " + " | ".join(c.replace("\n", " ").replace("|", "\\|").strip() for c in r) + " |"

    yield cells(header)
    yield "| " + " | ".join("---" for _ in range(width)) + " |"
    for row in iterator:
        yield cells(row)


def _render(lines: Iterable[str]) -> str:
    return "\n".join(lines).strip() + "\n"


def _pptx_lines(path: Path) -> Iterator[str]:
    from pptx import Presentation

    for i, slide in enumerate(Presentation(str(path)).slides, 1):
        yield f"## Slide {i}"
        for shape in slide.shapes:
            if shape.has_table:
                yield from _table_lines([[c.text.strip() for c in row.cells] for row in shape.table.rows])
            elif shape.has_text_frame and shape.text_frame.text.strip():
                yield shape.text_frame.text.strip()


def pptx_to_md(path: Path) -> str:
    return _render(_pptx_lines(path))


def _xlsx_lines(path: Path) -> Iterator[str]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=True, read_only=True)
    try:
        for ws in wb.worksheets:
            yield f"## Sheet: {ws.title}"
            rows = (
                ["" if cell is None else str(cell) for cell in row]
                for row in ws.iter_rows(values_only=True)
            )
            yield from _table_lines(row for row in rows if any(cell.strip() for cell in row))
    finally:
        wb.close()


def xlsx_to_md(path: Path) -> str:
    return _render(_xlsx_lines(path))


def _docx_lines(path: Path) -> Iterator[str]:
    from docx import Document

    doc = Document(str(path))
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name or ""
        if style == "Title":
            yield "# " + text
        elif style.startswith("Heading"):
            tail = style.split()[-1]
            yield "#" * min(int(tail) if tail.isdigit() else 2, 6) + " " + text
        else:
            yield text
    for table in doc.tables:
        yield from _table_lines([[c.text.strip() for c in row.cells] for row in table.rows])


def docx_to_md(path: Path) -> str:
    return _render(_docx_lines(path))


_CONVERTERS = {".pptx": pptx_to_md, ".xlsx": xlsx_to_md, ".docx": docx_to_md}
_LINE_CONVERTERS = {".pptx": _pptx_lines, ".xlsx": _xlsx_lines, ".docx": _docx_lines}


def convert_file(path: Path) -> str | None:
    """Render one office file to markdown, or None if it is not an office format."""
    fn = _CONVERTERS.get(path.suffix.lower())
    return fn(path) if fn else None


def _validate_archive_budget(path: Path) -> None:
    max_bytes = _positive_env("KBC_MAX_OFFICE_ARCHIVE_UNPACKED_BYTES", DEFAULT_MAX_ARCHIVE_UNPACKED_BYTES)
    max_files = _positive_env("KBC_MAX_OFFICE_ARCHIVE_FILES", DEFAULT_MAX_ARCHIVE_FILES)
    try:
        with zipfile.ZipFile(path) as archive:
            infos = archive.infolist()
            unpacked = sum(info.file_size for info in infos)
    except (OSError, zipfile.BadZipFile):
        # Preserve the existing per-file corrupt-input handling; the real parser
        # will report the format-specific error to convert_tree.
        return
    if len(infos) > max_files:
        raise OfficeIngestLimitExceeded(
            f"{path.name} contains {len(infos)} archive entries; limit is {max_files}"
        )
    if unpacked > max_bytes:
        raise OfficeIngestLimitExceeded(
            f"{path.name} expands to {unpacked} bytes; Office archive limit is {max_bytes}"
        )


def _write_bounded_markdown(path: Path, dest: Path, remaining: int) -> int:
    converter = _LINE_CONVERTERS[path.suffix.lower()]
    _validate_archive_budget(path)
    fd, temp_name = tempfile.mkstemp(prefix=f".{dest.name}.", suffix=".tmp", dir=dest.parent)
    written = 0
    lines = 0
    try:
        with os.fdopen(fd, "wb") as output:
            for line in converter(path):
                payload = (line + "\n").encode("utf-8")
                if written + len(payload) > remaining:
                    raise OfficeIngestLimitExceeded(
                        f"{path.name} derived markdown exceeds the remaining {remaining}-byte Office budget"
                    )
                output.write(payload)
                written += len(payload)
                lines += 1
            if lines == 0:
                # Preserve the previous converter contract for a valid but
                # empty Office file: it produced a one-newline readable
                # sidecar, rather than making the supported source look opaque.
                if remaining < 1:
                    raise OfficeIngestLimitExceeded(
                        f"{path.name} derived markdown exceeds the remaining {remaining}-byte Office budget"
                    )
                output.write(b"\n")
                written = 1
        os.replace(temp_name, dest)
        return written
    finally:
        if os.path.exists(temp_name):
            os.unlink(temp_name)


def convert_tree(root: str) -> tuple[list[tuple[str, str]], list[tuple[str, str]]]:
    """Walk `root`, render each .pptx/.xlsx/.docx to a sibling `<name>.md`.

    Returns (converted, errors): converted = [(src_rel, md_rel)], errors =
    [(src_rel, message)]. A corrupt parser input is recorded and skipped so one
    bad deck does not sink the KB install. Declared resource-budget violations
    fail the atomic snapshot commit instead of silently leaving a supported
    source unreadable. Idempotent: an existing `<name>.md` is left untouched (a
    re-install of the same bundle is a no-op)."""
    r = Path(root)
    converted: list[tuple[str, str]] = []
    errors: list[tuple[str, str]] = []
    max_derived = _positive_env("KBC_MAX_OFFICE_DERIVED_BYTES", DEFAULT_MAX_DERIVED_BYTES)
    derived = 0
    if not r.is_dir():
        return converted, errors
    for f in sorted(r.rglob("*")):
        if not f.is_file() or f.suffix.lower() not in OFFICE_EXTS:
            continue
        rel = f.relative_to(r).as_posix()
        dest = f.with_name(f.name + ".md")  # deck.pptx -> deck.pptx.md
        if dest.exists():
            continue
        try:
            written = _write_bounded_markdown(f, dest, max_derived - derived)
        except OfficeIngestLimitExceeded:
            # Resource-budget violations are not corrupt-file anomalies. The
            # source type is otherwise supported, so fail the atomic snapshot
            # commit rather than silently compiling an unreadable binary.
            raise
        except Exception as e:  # any parser failure is per-file non-fatal (fail-open boundary)
            errors.append((rel, repr(e)))
            continue
        if written == 0:
            continue
        derived += written
        converted.append((rel, dest.relative_to(r).as_posix()))
    return converted, errors
