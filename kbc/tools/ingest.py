#!/usr/bin/env python3
"""ingest.py — 异构本地文件 → 归一 markdown + 精确出处(域无关,机制层)。

把任意本地文档树解析成 OKF 编译可吸收的中间态:每段内容前带 <!-- @prov ... --> 出处锚,
出处(文件 + 页/幻灯片/工作表)同时写进 <name>.provenance.json。
出处是文件系统本地的 → "回源 = 重读本地文件",无需任何外部源接入。

引擎可插拔(按扩展名分派);v1 = 纯 python 轻量引擎:
  pdf  → pdfplumber(逐页文本 + 表)       pptx → python-pptx(逐幻灯片文本,图片登记)
  xlsx → openpyxl(逐表渲染 markdown 表)   图片 → 登记 + 留出处,语义交编译阶段视觉模型
  其它文本(md/txt/py/yaml/json/...) → 原样透传
高保真升级件(OCR / 扫描件 / 复杂版面)= Docling,后接(见 README)。
本工具不内置任何库的领域知识。
"""
import argparse
import json
import sys
from pathlib import Path

IMAGE_EXT = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"}
TEXT_EXT = {".md", ".markdown", ".txt", ".py", ".yaml", ".yml", ".json", ".toml",
            ".ini", ".cfg", ".sh", ".sql", ".csv", ".rst", ".log", ".xml", ".html"}
MAX_TABLE_ROWS = 200


class Doc:
    """累积 markdown 段 + 出处锚。"""

    def __init__(self, src, engine, fmt):
        self.src = src
        self.engine = engine
        self.fmt = fmt
        self.parts = []
        self.anchors = []

    def add(self, aid, loc, text):
        self.parts.append(f"<!-- @prov {aid} src={self.src.name} {loc} -->\n{text}\n\n")
        self.anchors.append({"id": aid, "loc": loc})

    def markdown(self):
        fm = (f"---\ntype: ingested\nsource: {self.src}\nengine: {self.engine}\n"
              f"format: {self.fmt}\nanchors: {len(self.anchors)}\n---\n\n")
        return fm + "".join(self.parts)

    def provenance(self):
        return {"source": str(self.src), "engine": self.engine,
                "format": self.fmt, "anchors": self.anchors}


def md_table(rows):
    rows = [r for r in rows if r is not None][:MAX_TABLE_ROWS]
    if not rows:
        return ""
    width = max(len(r) for r in rows)

    def fmt(r):
        cells = [("" if c is None else str(c)).replace("\n", " ").replace("|", "\\|") for c in r]
        cells += [""] * (width - len(cells))
        return "| " + " | ".join(cells) + " |"

    out = [fmt(rows[0]), "| " + " | ".join(["---"] * width) + " |"]
    out += [fmt(r) for r in rows[1:]]
    return "\n".join(out)


def parse_pdf(path):
    import pdfplumber
    d = Doc(path, "pdfplumber", "pdf")
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            body = (page.extract_text() or "").strip()
            for tbl in page.extract_tables() or []:
                body += "\n\n" + md_table(tbl)
            if not body.strip():
                body = "(本页无可抽取文本 — 可能是扫描/纯图页,待 OCR 或视觉描述)"
            d.add(f"p{i}", f"page={i}", body.strip())
    return d


def parse_pptx(path):
    from pptx import Presentation
    d = Doc(path, "python-pptx", "pptx")
    for i, slide in enumerate(Presentation(path).slides, 1):
        lines, imgs = [], 0
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())
            if shape.shape_type == 13:  # PICTURE
                imgs += 1
        body = "\n\n".join(lines)
        if imgs:
            body += f"\n\n> 🖼 本页含 {imgs} 张图片 — 待编译阶段视觉描述(出处 slide={i})"
        d.add(f"s{i}", f"slide={i}", body.strip() or "(空页)")
    return d


def parse_xlsx(path):
    from openpyxl import load_workbook
    d = Doc(path, "openpyxl", "xlsx")
    wb = load_workbook(path, data_only=True, read_only=True)
    for i, ws in enumerate(wb.worksheets, 1):
        rows = list(ws.iter_rows(values_only=True))
        ncols = max((len(r) for r in rows), default=0)
        body = f"### {ws.title}\n\n" + (md_table(rows) if rows else "(空表)")
        d.add(f"sh{i}", f"sheet={ws.title} rows={len(rows)} cols={ncols}", body)
    wb.close()
    return d


def parse_image(path):
    from PIL import Image
    d = Doc(path, "pillow", "image")
    with Image.open(path) as im:
        w, h = im.size
    d.add("img1", f"region=full wh={w}x{h}",
          "> 🖼 图片待视觉描述(交编译阶段视觉模型 / 多模态协议);ingest 仅登记 + 留出处。")
    return d


def parse_text(path):
    d = Doc(path, "passthrough", path.suffix.lstrip(".") or "text")
    d.add("t1", "whole", path.read_text(encoding="utf-8", errors="replace").strip())
    return d


PARSERS = {".pdf": parse_pdf, ".pptx": parse_pptx, ".xlsx": parse_xlsx}


def parse_file(path):
    ext = path.suffix.lower()
    if ext in PARSERS:
        return PARSERS[ext](path)
    if ext in IMAGE_EXT:
        return parse_image(path)
    if ext in TEXT_EXT:
        return parse_text(path)
    return None


def safe_name(path, scan_root):
    rel = path.relative_to(scan_root).as_posix() if scan_root else path.name
    return rel.replace("/", "__").replace(" ", "_")


def main():
    ap = argparse.ArgumentParser(description="异构文件 → 归一 markdown + 出处(域无关)")
    ap.add_argument("--src", required=True, help="文件或目录")
    ap.add_argument("--out", required=True, help="输出目录")
    a = ap.parse_args()

    src = Path(a.src).expanduser()
    out = Path(a.out).expanduser()
    out.mkdir(parents=True, exist_ok=True)
    if src.is_file():
        files, scan_root = [src], src.parent
    elif src.is_dir():
        files, scan_root = sorted(p for p in src.rglob("*") if p.is_file()), src
    else:
        sys.exit(f"✗ 源不存在: {src}")

    ok = skip = fail = 0
    for f in files:
        try:
            doc = parse_file(f)
        except Exception as e:                       # 批处理边界:单文件失败不拖垮整批
            print(f"  ✗ {f.name}: {type(e).__name__}: {e}")
            fail += 1
            continue
        if doc is None:
            skip += 1
            continue
        stem = safe_name(f, scan_root)
        (out / f"{stem}.md").write_text(doc.markdown(), encoding="utf-8")
        (out / f"{stem}.provenance.json").write_text(
            json.dumps(doc.provenance(), ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"  ✅ {f.name}  [{doc.engine}] {len(doc.anchors)} 锚")
        ok += 1

    print(f"\n解析 {ok} · 跳过(不支持) {skip} · 失败 {fail}")
    sys.exit(1 if fail else 0)


if __name__ == "__main__":
    main()
